"""
Q-TRON Trading Console — Design Concept Proposal
Premium dark theme presentation generator
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# ── Color Palette ──
BG_MAIN    = RGBColor(0x06, 0x08, 0x10)
BG_CARD    = RGBColor(0x0E, 0x11, 0x18)
BG_HOVER   = RGBColor(0x16, 0x1A, 0x24)
BG_ACTIVE  = RGBColor(0x1E, 0x22, 0x30)
GREEN      = RGBColor(0x00, 0xFF, 0x88)
RED        = RGBColor(0xFF, 0x33, 0x44)
YELLOW     = RGBColor(0xFF, 0xDD, 0x00)
BLUE       = RGBColor(0x66, 0x99, 0xFF)
CYAN       = RGBColor(0x44, 0xCC, 0xCC)
WHITE      = RGBColor(0xF0, 0xF0, 0xF0)
DIM        = RGBColor(0x60, 0x68, 0x80)
DIMMER     = RGBColor(0x40, 0x48, 0x58)
ACCENT_GRN = RGBColor(0x00, 0xCC, 0x6A)
ACCENT_BLU = RGBColor(0x33, 0x66, 0xCC)
ORANGE     = RGBColor(0xFF, 0x88, 0x00)

# ── Helpers ──
def set_slide_bg(slide, color=BG_MAIN):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill_color, border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, x, y, w, h, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font_name="Segoe UI", valign=MSO_ANCHOR.TOP,
             line_spacing=None):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    if valign:
        tf.auto_size = None
        from pptx.oxml.ns import qn
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        bodyPr.set('anchor', {MSO_ANCHOR.TOP: 't', MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b'}[valign])
    if line_spacing:
        p.line_spacing = Pt(line_spacing)
    return txBox

def add_multiline(slide, lines, x, y, w, h, size=12, color=WHITE,
                  font_name="Segoe UI", line_spacing=18, bold=False, align=PP_ALIGN.LEFT):
    """lines: list of (text, color, bold, size) or just str"""
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, clr, b, sz = item, color, bold, size
        else:
            txt = item[0]
            clr = item[1] if len(item) > 1 else color
            b = item[2] if len(item) > 2 else bold
            sz = item[3] if len(item) > 3 else size
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = b
        p.font.name = font_name
        p.alignment = align
        p.space_before = Pt(2)
        p.space_after = Pt(2)
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
    return txBox

def add_circle(slide, x, y, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_line(slide, x, y, w, color=DIMMER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_kpi_card(slide, x, y, w, h, label, value, sub="", value_color=GREEN):
    add_rounded_rect(slide, x, y, w, h, BG_CARD, border_color=DIMMER)
    add_text(slide, label, x + Inches(0.15), y + Inches(0.08), w - Inches(0.3), Inches(0.25),
             size=9, color=DIM, font_name="Consolas")
    add_text(slide, value, x + Inches(0.15), y + Inches(0.28), w - Inches(0.3), Inches(0.35),
             size=20, color=value_color, bold=True, font_name="Consolas")
    if sub:
        add_text(slide, sub, x + Inches(0.15), y + h - Inches(0.22), w - Inches(0.3), Inches(0.2),
                 size=8, color=DIM, font_name="Consolas")

def add_status_led(slide, x, y, label, color_led):
    add_circle(slide, x, y + Inches(0.02), Inches(0.12), color_led)
    add_text(slide, label, x + Inches(0.18), y - Inches(0.02), Inches(0.8), Inches(0.2),
             size=8, color=DIM, font_name="Consolas", bold=True)

# ── Presentation ──
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = Inches(13.333)
SH = Inches(7.5)


# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)

# Top accent line
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.04), GREEN)

# Left vertical accent
add_rect(slide, Inches(0.6), Inches(1.5), Inches(0.04), Inches(3.5), GREEN)

# Main title
add_text(slide, "Q-TRON", Inches(1.0), Inches(1.8), Inches(8), Inches(1.2),
         size=64, color=WHITE, bold=True, font_name="Segoe UI Light")
add_text(slide, "TRADING CONSOLE", Inches(1.0), Inches(2.8), Inches(8), Inches(0.8),
         size=40, color=GREEN, bold=True, font_name="Consolas")

# Subtitle
add_text(slide, "Next-Gen Quantitative Trading GUI", Inches(1.0), Inches(3.8), Inches(6), Inches(0.5),
         size=18, color=DIM, font_name="Segoe UI")
add_text(slide, "Design Concept Proposal", Inches(1.0), Inches(4.2), Inches(6), Inches(0.5),
         size=18, color=DIM, font_name="Segoe UI")

# Bottom info bar
add_rect(slide, Inches(0), SH - Inches(0.8), SW, Inches(0.8), BG_CARD)
add_text(slide, "DESIGN STUDIO X", Inches(1.0), SH - Inches(0.65), Inches(3), Inches(0.4),
         size=12, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "Q-TRON QUANT", Inches(4.5), SH - Inches(0.65), Inches(3), Inches(0.4),
         size=12, color=BLUE, bold=True, font_name="Consolas")
add_text(slide, "2026.04", Inches(10.5), SH - Inches(0.65), Inches(2), Inches(0.4),
         size=12, color=DIM, font_name="Consolas", align=PP_ALIGN.RIGHT)

# Right side decorative elements - circuit-like pattern
for i in range(6):
    y_pos = Inches(1.0 + i * 0.9)
    add_rect(slide, Inches(10.5), y_pos, Inches(2.0), Pt(1), DIMMER)
    add_circle(slide, Inches(12.3), y_pos - Inches(0.03), Inches(0.06), DIMMER)


# ============================================================
# SLIDE 2: EXECUTIVE SUMMARY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Title area
add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "01", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "PROJECT OVERVIEW", Inches(1.5), Inches(0.15), Inches(6), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uD504\uB85C\uC81D\uD2B8 \uAC1C\uC694", Inches(1.5), Inches(0.6), Inches(6), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Hero quote
add_rounded_rect(slide, Inches(0.6), Inches(1.5), Inches(12.1), Inches(1.2), BG_HOVER, border_color=GREEN)
add_rect(slide, Inches(0.6), Inches(1.5), Inches(0.06), Inches(1.2), GREEN)
add_text(slide, "\"1\uCD08 \uC548\uC5D0 \uC2DC\uC7A5 \uC0C1\uD0DC\uB97C \uD30C\uC545\uD558\uACE0, 3\uCD08 \uC548\uC5D0 \uB9AC\uC2A4\uD06C\uB97C \uD310\uB2E8\uD55C\uB2E4\"",
         Inches(1.0), Inches(1.65), Inches(11.5), Inches(0.8),
         size=22, color=GREEN, bold=True, font_name="Segoe UI", align=PP_ALIGN.CENTER)

# 3 column cards
card_w = Inches(3.7)
card_h = Inches(3.6)
card_y = Inches(3.1)
topics = [
    ("\uBBF8\uC158", "Mission", "\uAE30\uC874 PyQt5 Dock \uAE30\uBC18 \uBAA8\uB2C8\uD130\uB97C\n\uCC28\uC138\uB300 \uD2B8\uB808\uC774\uB529 \uCF58\uC194\uB85C \uC7AC\uC124\uACC4",
     ["\uC2E4\uC2DC\uAC04 \uC758\uC0AC\uACB0\uC815 \uC9C0\uC6D0", "\uB9AC\uC2A4\uD06C \uAC00\uC2DC\uD654 \uADF9\uB300\uD654", "\uC6B4\uC601 \uC2E0\uB8B0\uC131 \uD655\uBCF4"],
     GREEN),
    ("\uBC94\uC704", "Scope", "20\uC885\uBAA9 \uB85C\uD14C\uC774\uC158 \uC804\uB7B5\n\uC6D4\uAC04 \uB9AC\uBC38\uB7F0\uC2F1 + Trail Stop",
     ["BuyPermission \uC0C1\uD0DC\uAE30\uACC4", "AI Advisor 4-Phase \uD1B5\uD569", "DD Guard \uB9AC\uC2A4\uD06C \uAD00\uB9AC"],
     BLUE),
    ("\uBAA9\uD45C", "Goal", "\uD5C8\uB2E4\uD55C \uC815\uBCF4 \uB098\uC5F4\uC744 \uBC97\uC5B4\uB098\n\uC6B0\uC120\uC21C\uC704 \uAE30\uBC18 \uC2DC\uAC01 \uACC4\uCE35\uAD6C\uC870",
     ["\uAD00\uC81C \uC13C\uD130 \uD328\uB7EC\uB2E4\uC784", "\uAE00\uB79C\uC2A4 \uAC00\uB2A5\uD55C \uB300\uC2DC\uBCF4\uB4DC", "\uB9AC\uC2A4\uD06C \uC990\uAC01 \uC2DC\uAC01\uD654"],
     CYAN),
]

for i, (ko_label, en_label, desc, bullets, accent) in enumerate(topics):
    cx = Inches(0.6 + i * 4.1)
    add_rounded_rect(slide, cx, card_y, card_w, card_h, BG_CARD, border_color=DIMMER)
    # Accent top bar
    add_rect(slide, cx, card_y, card_w, Inches(0.04), accent)
    # Number
    add_text(slide, f"0{i+1}", cx + Inches(0.2), card_y + Inches(0.2), Inches(0.6), Inches(0.4),
             size=24, color=accent, bold=True, font_name="Consolas")
    # Title
    add_text(slide, en_label, cx + Inches(0.8), card_y + Inches(0.2), Inches(2.5), Inches(0.35),
             size=18, color=WHITE, bold=True)
    add_text(slide, ko_label, cx + Inches(0.8), card_y + Inches(0.5), Inches(2.5), Inches(0.3),
             size=11, color=DIM)
    # Desc
    add_text(slide, desc, cx + Inches(0.2), card_y + Inches(0.9), card_w - Inches(0.4), Inches(0.8),
             size=11, color=WHITE, line_spacing=16)
    # Bullets
    bullet_lines = [(f"\u25B8  {b}", accent, False, 10) for b in bullets]
    add_multiline(slide, bullet_lines, cx + Inches(0.2), card_y + Inches(1.9),
                  card_w - Inches(0.4), Inches(1.5), line_spacing=16)


# ============================================================
# SLIDE 3: AS-IS ANALYSIS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "02", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "AS-IS ANALYSIS", Inches(1.5), Inches(0.15), Inches(6), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uD604\uC7AC \uC2DC\uC2A4\uD15C \uBD84\uC11D", Inches(1.5), Inches(0.6), Inches(6), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Current layout diagram - 4 rows
row_data = [
    ("Row 1  KPI Strip", "Decision Card + Today PnL + MDD + Positions + Cash + Cumulative", GREEN),
    ("Row 2  Charts", "PnL Bar Chart (60%) + Equity vs KOSPI Line (40%)", BLUE),
    ("Row 3  Portfolio", "Sector Donut + Positions Table + Advisor Alerts", CYAN),
    ("Row 4  Activity", "Recent Trades (60%) + System Status Badges (40%)", YELLOW),
]

diagram_x = Inches(0.6)
diagram_y = Inches(1.4)
for i, (label, desc, accent) in enumerate(row_data):
    ry = diagram_y + Inches(i * 0.75)
    add_rounded_rect(slide, diagram_x, ry, Inches(5.8), Inches(0.6), BG_CARD, border_color=DIMMER)
    add_rect(slide, diagram_x, ry, Inches(0.05), Inches(0.6), accent)
    add_text(slide, label, diagram_x + Inches(0.2), ry + Inches(0.03), Inches(2.2), Inches(0.25),
             size=10, color=accent, bold=True, font_name="Consolas")
    add_text(slide, desc, diagram_x + Inches(0.2), ry + Inches(0.28), Inches(5.3), Inches(0.25),
             size=9, color=DIM, font_name="Consolas")

# Strengths & Weaknesses
sw_x = Inches(7.0)
add_text(slide, "\u2714  STRENGTHS", sw_x, Inches(1.4), Inches(5.5), Inches(0.35),
         size=14, color=GREEN, bold=True, font_name="Consolas")
strengths = [
    "\u25B8  \uD48D\uBD80\uD55C \uB370\uC774\uD130 \uD45C\uC2DC (KPI, \uCC28\uD2B8, \uD14C\uC774\uBE14)",
    "\u25B8  \uB2E4\uD06C \uD14C\uB9C8 + Glassmorphism \uC801\uC6A9",
    "\u25B8  Dock \uC7AC\uBC30\uCE58 \uAC00\uB2A5 (\uD30C\uC6CC\uC720\uC800)",
    "\u25B8  3\uCD08 \uC8FC\uAE30 \uC2E4\uC2DC\uAC04 \uAC31\uC2E0",
]
s_lines = [(s, WHITE, False, 10) for s in strengths]
add_multiline(slide, s_lines, sw_x, Inches(1.8), Inches(5.5), Inches(1.2), line_spacing=16)

add_text(slide, "\u2716  LIMITATIONS", sw_x, Inches(3.2), Inches(5.5), Inches(0.35),
         size=14, color=RED, bold=True, font_name="Consolas")
limits = [
    ("\u25B8  \uC815\uBCF4 \uC6B0\uC120\uC21C\uC704 \uBD88\uBA85\uD655 \u2014 \uBAA8\uB4E0 \uB370\uC774\uD130\uAC00 \uB3D9\uC77C \uBE44\uC911", RED),
    ("\u25B8  \uC2DC\uAC01\uC801 \uACC4\uCE35\uAD6C\uC870 \uBD80\uC7AC \u2014 KPI\uC640 \uCC28\uD2B8 \uAD6C\uBD84 \uC57D\uD568", RED),
    ("\u25B8  \uC54C\uB9BC \uD1B5\uD569 \uBBF8\uD761 \u2014 Advisor \uACB0\uACFC\uAC00 \uC791\uC740 \uCE74\uB4DC\uC5D0 \uBB3B\uD600 \uC788\uC74C", RED),
    ("\u25B8  \uB9AC\uC2A4\uD06C \uC2DC\uAC01\uD654 \uBD80\uC871 \u2014 Trail Gap, DD Guard \uC9C1\uAD00\uC131 \uBD80\uC871", RED),
]
l_lines = [(t, c, False, 10) for t, c in limits]
add_multiline(slide, l_lines, sw_x, Inches(3.6), Inches(5.5), Inches(1.5), line_spacing=16)

# Stats
add_rect(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(0.04), DIMMER)
stats = [
    ("8", "QDockWidgets"),
    ("4", "Row Layout"),
    ("5", "KPI Cards"),
    ("2", "Chart Types"),
    ("5", "Status Badges"),
    ("3s", "Refresh Cycle"),
]
for i, (val, label) in enumerate(stats):
    sx = Inches(0.6 + i * 2.05)
    add_text(slide, val, sx, Inches(5.5), Inches(1.5), Inches(0.5),
             size=28, color=GREEN, bold=True, font_name="Consolas")
    add_text(slide, label, sx, Inches(6.0), Inches(1.5), Inches(0.3),
             size=10, color=DIM, font_name="Consolas")


# ============================================================
# SLIDE 4: DESIGN PHILOSOPHY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "03", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "DESIGN PHILOSOPHY", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "Command Center Paradigm", Inches(1.5), Inches(0.6), Inches(6), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# 3 principles - large cards
principles = [
    ("GLANCEABLE", "\uD55C\uB208\uC5D0 \uD30C\uC545", "1\uCD08 \uB8F0",
     "\uC804\uCCB4 \uD3EC\uD2B8\uD3F4\uB9AC\uC624 \uC0C1\uD0DC\uB97C\n\uB2E8 1\uCD08\uB9CC\uC5D0 \uD30C\uC545 \uAC00\uB2A5\uD55C\n\uC2DC\uAC01 \uACC4\uCE35\uAD6C\uC870 \uC124\uACC4",
     GREEN, "\uD56D\uACF5\uAD00\uC81C \uB300\uC2DC\uBCF4\uB4DC"),
    ("ACTIONABLE", "\uC989\uC2DC \uD310\uB2E8", "3\uCD08 \uB8F0",
     "\uB9AC\uC2A4\uD06C \uC0C1\uD669 \uBC1C\uC0DD \uC2DC\n3\uCD08 \uC774\uB0B4 \uC704\uD5D8 \uC694\uC778 \uC2DD\uBCC4\uACFC\n\uB300\uC751 \uBC29\uD5A5 \uD310\uB2E8 \uAC00\uB2A5",
     YELLOW, "\uAE08\uC735 \uD130\uBBF8\uB110 UX"),
    ("TRUSTWORTHY", "\uC2E0\uB8B0 \uAC00\uB2A5", "\uD56D\uC2DC \uAC00\uC2DC",
     "\uB370\uC774\uD130 \uC2E0\uC120\uB3C4, \uC2DC\uC2A4\uD15C \uC0C1\uD0DC,\n\uBE0C\uB85C\uCEE4 \uC5F0\uACB0 \uC0C1\uD0DC\uAC00\n\uD56D\uC0C1 \uB208\uC5D0 \uBCF4\uC774\uB294 \uC124\uACC4",
     CYAN, "NASA Mission Control"),
]

for i, (title, ko, rule, desc, accent, inspo) in enumerate(principles):
    cx = Inches(0.6 + i * 4.1)
    cy = Inches(1.5)
    cw = Inches(3.8)
    ch = Inches(4.5)
    add_rounded_rect(slide, cx, cy, cw, ch, BG_CARD, border_color=DIMMER)

    # Large number circle
    circle = add_circle(slide, cx + Inches(0.2), cy + Inches(0.3), Inches(0.6), accent)
    add_text(slide, str(i+1), cx + Inches(0.2), cy + Inches(0.32), Inches(0.6), Inches(0.5),
             size=22, color=BG_MAIN, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)

    add_text(slide, title, cx + Inches(1.0), cy + Inches(0.3), Inches(2.6), Inches(0.35),
             size=18, color=accent, bold=True, font_name="Consolas")
    add_text(slide, ko, cx + Inches(1.0), cy + Inches(0.65), Inches(2.6), Inches(0.3),
             size=12, color=WHITE)

    # Rule badge
    add_rounded_rect(slide, cx + Inches(0.2), cy + Inches(1.1), Inches(1.5), Inches(0.35), BG_HOVER, border_color=accent)
    add_text(slide, rule, cx + Inches(0.2), cy + Inches(1.13), Inches(1.5), Inches(0.3),
             size=11, color=accent, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)

    # Description
    add_text(slide, desc, cx + Inches(0.2), cy + Inches(1.7), cw - Inches(0.4), Inches(1.2),
             size=11, color=WHITE, line_spacing=17)

    # Inspiration
    add_rect(slide, cx + Inches(0.2), cy + Inches(3.5), cw - Inches(0.4), Pt(1), DIMMER)
    add_text(slide, f"Inspired by: {inspo}", cx + Inches(0.2), cy + Inches(3.7), cw - Inches(0.4), Inches(0.3),
             size=9, color=DIM, font_name="Consolas")

# Bottom tagline
add_text(slide, "INSPIRATION:  \uD56D\uACF5\uAD00\uC81C  \u00B7  \uAE08\uC735 \uD130\uBBF8\uB110  \u00B7  NASA Mission Control",
         Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         size=12, color=DIM, font_name="Consolas", align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: INFORMATION ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "04", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "INFORMATION ARCHITECTURE", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uC815\uBCF4 \uC544\uD0A4\uD14D\uCC98  \u2014  3-Tier Priority System", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# 3 tiers - pyramid-like
tiers = [
    ("TIER 1", "ALWAYS VISIBLE", GREEN,
     ["BuyPermission \uC0C1\uD0DC", "Today PnL %", "Current MDD %", "\uD3EC\uC9C0\uC158 \uC218 (N/20)"],
     "\uD56D\uC0C1 \uD654\uBA74\uC5D0 \uB178\uCD9C \u2014 \uAC00\uC7A5 \uB113\uC740 \uBA74\uC801, \uAC00\uC7A5 \uD070 \uAE00\uC790"),
    ("TIER 2", "GLANCEABLE", YELLOW,
     ["Equity Curve \uCC28\uD2B8", "\uC139\uD130\uBCC4 \uBE44\uC911", "Trail Stop \uADFC\uC811 \uC885\uBAA9", "\uB9AC\uC2A4\uD06C \uC218\uC900 \uAC8C\uC774\uC9C0"],
     "\uC2DC\uC120\uC744 \uB3CC\uB9AC\uBA74 \uD655\uC778 \uAC00\uB2A5 \u2014 \uC911\uAC04 \uBA74\uC801, \uC2DC\uAC01\uC801 \uCC28\uD2B8 \uC911\uC2EC"),
    ("TIER 3", "ON-DEMAND", BLUE,
     ["\uAC70\uB798 \uB0B4\uC5ED \uC0C1\uC138", "\uC2DC\uC2A4\uD15C \uB85C\uADF8", "AI Advisor \uCD94\uCC9C", "\uD30C\uB77C\uBBF8\uD130 \uC870\uC815 \uC81C\uC548"],
     "\uD074\uB9AD/\uD38C\uCE68\uC73C\uB85C \uD655\uC778 \u2014 \uC0C1\uC138 \uD14D\uC2A4\uD2B8, \uD14C\uC774\uBE14, \uD328\uB110"),
]

for i, (tier, name, accent, items, desc) in enumerate(tiers):
    ty = Inches(1.4 + i * 1.85)
    # Width narrows from tier 1 to 3 (pyramid effect)
    tw = Inches(12.1 - i * 1.5)
    tx = Inches(0.6 + i * 0.75)

    add_rounded_rect(slide, tx, ty, tw, Inches(1.6), BG_CARD, border_color=accent)
    add_rect(slide, tx, ty, tw, Inches(0.04), accent)

    # Tier label
    add_text(slide, tier, tx + Inches(0.2), ty + Inches(0.15), Inches(1.2), Inches(0.3),
             size=14, color=accent, bold=True, font_name="Consolas")
    add_text(slide, name, tx + Inches(1.4), ty + Inches(0.15), Inches(2.5), Inches(0.3),
             size=14, color=WHITE, bold=True)

    # Items in a row
    iw = (tw - Inches(0.4)) / len(items)
    for j, item in enumerate(items):
        ix = tx + Inches(0.2) + iw * j
        add_rounded_rect(slide, ix, ty + Inches(0.55), iw - Inches(0.1), Inches(0.4), BG_HOVER)
        add_text(slide, item, ix + Inches(0.05), ty + Inches(0.58), iw - Inches(0.2), Inches(0.35),
                 size=9, color=WHITE, font_name="Consolas", align=PP_ALIGN.CENTER)

    # Description
    add_text(slide, desc, tx + Inches(0.2), ty + Inches(1.1), tw - Inches(0.4), Inches(0.3),
             size=10, color=DIM, font_name="Segoe UI")


# ============================================================
# SLIDE 6: LAYOUT CONCEPT - MISSION CONTROL
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "05", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "LAYOUT CONCEPT", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "Mission Control \uB808\uC774\uC544\uC6C3", Inches(1.5), Inches(0.6), Inches(6), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Wireframe mock
wire_x = Inches(0.8)
wire_y = Inches(1.4)
wire_w = Inches(11.7)
wire_h = Inches(5.5)

# Outer frame
add_rounded_rect(slide, wire_x, wire_y, wire_w, wire_h, BG_CARD, border_color=DIMMER)

# TOP BAR - Heartbeat
hb_y = wire_y + Inches(0.1)
add_rect(slide, wire_x + Inches(0.1), hb_y, wire_w - Inches(0.2), Inches(0.35), BG_HOVER)
labels = ["ENGINE", "STALE", "RISK", "REGIME", "UPDATED"]
colors = [GREEN, GREEN, YELLOW, BLUE, DIM]
for i, (lbl, clr) in enumerate(zip(labels, colors)):
    lx = wire_x + Inches(0.3 + i * 2.3)
    add_circle(slide, lx, hb_y + Inches(0.1), Inches(0.12), clr)
    add_text(slide, lbl, lx + Inches(0.18), hb_y + Inches(0.05), Inches(1.2), Inches(0.25),
             size=8, color=DIM, font_name="Consolas", bold=True)

# System Heartbeat Strip label
add_text(slide, "SYSTEM HEARTBEAT STRIP", wire_x + Inches(0.1), hb_y - Inches(0.01),
         wire_w - Inches(0.2), Inches(0.35),
         size=7, color=DIMMER, font_name="Consolas", align=PP_ALIGN.RIGHT)

# LEFT PANEL - Decision Hub (20%)
lp_x = wire_x + Inches(0.1)
lp_y = wire_y + Inches(0.55)
lp_w = Inches(2.2)
lp_h = Inches(4.8)
add_rounded_rect(slide, lp_x, lp_y, lp_w, lp_h, BG_HOVER, border_color=GREEN)
add_text(slide, "DECISION", lp_x + Inches(0.1), lp_y + Inches(0.1), lp_w - Inches(0.2), Inches(0.3),
         size=10, color=GREEN, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)
add_text(slide, "HUB", lp_x + Inches(0.1), lp_y + Inches(0.35), lp_w - Inches(0.2), Inches(0.25),
         size=10, color=GREEN, font_name="Consolas", align=PP_ALIGN.CENTER)
# Decision display
add_rounded_rect(slide, lp_x + Inches(0.15), lp_y + Inches(0.7), lp_w - Inches(0.3), Inches(0.6), RGBColor(0x00, 0x33, 0x22), border_color=GREEN)
add_text(slide, "BUY OK", lp_x + Inches(0.15), lp_y + Inches(0.78), lp_w - Inches(0.3), Inches(0.4),
         size=14, color=GREEN, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)
# KPI mini cards
kpi_labels = ["PnL", "MDD", "N/20", "Cash", "Cum"]
for j, kl in enumerate(kpi_labels):
    ky = lp_y + Inches(1.5 + j * 0.6)
    add_rounded_rect(slide, lp_x + Inches(0.15), ky, lp_w - Inches(0.3), Inches(0.45), BG_ACTIVE)
    add_text(slide, kl, lp_x + Inches(0.25), ky + Inches(0.05), Inches(0.8), Inches(0.18),
             size=7, color=DIM, font_name="Consolas")
    add_text(slide, "+2.3%", lp_x + Inches(0.25), ky + Inches(0.2), Inches(1.2), Inches(0.2),
             size=10, color=GREEN, bold=True, font_name="Consolas")

add_text(slide, "20%", lp_x, lp_y + lp_h + Inches(0.02), lp_w, Inches(0.2),
         size=8, color=DIM, font_name="Consolas", align=PP_ALIGN.CENTER)

# CENTER - Hero Chart (55%)
cp_x = lp_x + lp_w + Inches(0.1)
cp_y = lp_y
cp_w = Inches(6.5)
cp_h = Inches(3.5)
add_rounded_rect(slide, cp_x, cp_y, cp_w, cp_h, BG_HOVER, border_color=BLUE)
# Tab buttons
tabs = ["Equity Curve", "PnL Waterfall", "Risk Heatmap"]
for j, tab in enumerate(tabs):
    tx_tab = cp_x + Inches(0.1 + j * 2.1)
    tab_color = BLUE if j == 0 else DIMMER
    add_rounded_rect(slide, tx_tab, cp_y + Inches(0.1), Inches(1.9), Inches(0.3), BG_ACTIVE if j == 0 else BG_CARD)
    add_text(slide, tab, tx_tab, cp_y + Inches(0.12), Inches(1.9), Inches(0.25),
             size=8, color=tab_color, font_name="Consolas", align=PP_ALIGN.CENTER, bold=j==0)

# Mock chart lines
for j in range(5):
    ly = cp_y + Inches(0.7 + j * 0.55)
    add_rect(slide, cp_x + Inches(0.3), ly, cp_w - Inches(0.6), Pt(0.5), DIMMER)

add_text(slide, "HERO CHART AREA", cp_x, cp_y + Inches(1.4), cp_w, Inches(0.3),
         size=12, color=BLUE, font_name="Consolas", align=PP_ALIGN.CENTER, bold=True)
add_text(slide, "55%", cp_x, cp_y + cp_h + Inches(0.02), cp_w, Inches(0.2),
         size=8, color=DIM, font_name="Consolas", align=PP_ALIGN.CENTER)

# RIGHT PANEL - Alerts (25%)
rp_x = cp_x + cp_w + Inches(0.1)
rp_y = lp_y
rp_w = Inches(2.7)
rp_h = Inches(4.8)
add_rounded_rect(slide, rp_x, rp_y, rp_w, rp_h, BG_HOVER, border_color=ORANGE)
add_text(slide, "ALERT STREAM", rp_x + Inches(0.1), rp_y + Inches(0.1), rp_w - Inches(0.2), Inches(0.25),
         size=10, color=ORANGE, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)
# Mock alert cards
alert_colors = [RED, ORANGE, YELLOW, DIM]
alert_labels = ["CRITICAL", "HIGH", "MEDIUM", "LOW"]
for j in range(4):
    ay = rp_y + Inches(0.5 + j * 0.55)
    add_rounded_rect(slide, rp_x + Inches(0.1), ay, rp_w - Inches(0.2), Inches(0.42), BG_ACTIVE)
    add_rect(slide, rp_x + Inches(0.1), ay, Inches(0.04), Inches(0.42), alert_colors[j])
    add_text(slide, alert_labels[j], rp_x + Inches(0.25), ay + Inches(0.05), Inches(1.5), Inches(0.18),
             size=7, color=alert_colors[j], font_name="Consolas", bold=True)
    add_text(slide, "Alert message...", rp_x + Inches(0.25), ay + Inches(0.2), Inches(2.0), Inches(0.18),
             size=7, color=DIM, font_name="Consolas")

add_text(slide, "25%", rp_x, rp_y + rp_h + Inches(0.02), rp_w, Inches(0.2),
         size=8, color=DIM, font_name="Consolas", align=PP_ALIGN.CENTER)

# BOTTOM - Position Grid
bg_y = cp_y + cp_h + Inches(0.1)
bg_w = cp_w
bg_h = Inches(1.2)
add_rounded_rect(slide, cp_x, bg_y, bg_w, bg_h, BG_HOVER, border_color=CYAN)
add_text(slide, "POSITION GRID  \u2014  20 Stock Tiles", cp_x + Inches(0.2), bg_y + Inches(0.05), bg_w - Inches(0.4), Inches(0.25),
         size=9, color=CYAN, font_name="Consolas", bold=True)
# Mini tiles
for j in range(10):
    tx_tile = cp_x + Inches(0.1 + j * 0.63)
    tile_clr = GREEN if j % 3 != 0 else RED
    add_rounded_rect(slide, tx_tile, bg_y + Inches(0.35), Inches(0.55), Inches(0.7), BG_ACTIVE)
    add_rect(slide, tx_tile, bg_y + Inches(0.35), Inches(0.55), Inches(0.04), tile_clr)


# ============================================================
# SLIDE 7: DECISION HUB
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "06", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "DECISION HUB", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uD575\uC2EC \uCEF4\uD3EC\uB10C\uD2B8 #1  \u2014  \uB9E4\uC218 \uD5C8\uAC00 \uC0C1\uD0DC \uC911\uC2EC", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Three states side by side
states = [
    ("BUY OK", "\uC815\uC0C1 \uB9E4\uC218 \uAC00\uB2A5", GREEN, RGBColor(0x00, 0x33, 0x22),
     "Risk: NORMAL\nData: Fresh (<90s)\nRECON: Clean"),
    ("BUY LIMITED", "\uB9E4\uC218 \uCD95\uC18C (50~70%)", YELLOW, RGBColor(0x33, 0x2B, 0x00),
     "Risk: DD_CAUTION\nData: Warning (90~180s)\nRECON: Clean"),
    ("BUY BLOCKED", "\uB9E4\uC218 \uCC28\uB2E8", RED, RGBColor(0x33, 0x0A, 0x0E),
     "Risk: DD_CRITICAL\nData: STALE (>180s)\nRECON: Unreliable"),
]

for i, (label, ko, accent, bg_clr, details) in enumerate(states):
    cx = Inches(0.6 + i * 4.1)
    cy = Inches(1.5)
    cw = Inches(3.8)
    ch = Inches(2.8)
    add_rounded_rect(slide, cx, cy, cw, ch, bg_clr, border_color=accent)

    # Glow top
    add_rect(slide, cx, cy, cw, Inches(0.06), accent)

    # Status text
    add_text(slide, label, cx + Inches(0.2), cy + Inches(0.3), cw - Inches(0.4), Inches(0.6),
             size=28, color=accent, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)
    add_text(slide, ko, cx + Inches(0.2), cy + Inches(0.9), cw - Inches(0.4), Inches(0.3),
             size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Details
    for j, line in enumerate(details.split("\n")):
        add_text(slide, line, cx + Inches(0.3), cy + Inches(1.4 + j * 0.35), cw - Inches(0.6), Inches(0.3),
                 size=10, color=DIM, font_name="Consolas")

# KPI Stack description
add_rect(slide, Inches(0.6), Inches(4.6), Inches(12.1), Pt(1), DIMMER)
add_text(slide, "KPI VERTICAL STACK", Inches(0.6), Inches(4.8), Inches(3), Inches(0.3),
         size=14, color=WHITE, bold=True)
add_text(slide, "Decision Hub \uD558\uB2E8\uC5D0 5\uAC1C KPI \uCE74\uB4DC \uC218\uC9C1 \uBC30\uCE58", Inches(0.6), Inches(5.1), Inches(5), Inches(0.3),
         size=11, color=DIM)

# 5 KPI mini cards in a row
kpi_data = [
    ("TODAY PNL", "+1.82%", GREEN),
    ("CURRENT MDD", "-3.2%", YELLOW),
    ("POSITIONS", "18 / 20", WHITE),
    ("CASH RATIO", "12.4%", CYAN),
    ("CUMULATIVE", "+47.3%", GREEN),
]
for i, (kpi_label, kpi_val, kpi_clr) in enumerate(kpi_data):
    kx = Inches(0.6 + i * 2.45)
    ky = Inches(5.6)
    add_rounded_rect(slide, kx, ky, Inches(2.25), Inches(1.1), BG_CARD, border_color=DIMMER)
    add_text(slide, kpi_label, kx + Inches(0.15), ky + Inches(0.1), Inches(1.9), Inches(0.2),
             size=8, color=DIM, font_name="Consolas")
    add_text(slide, kpi_val, kx + Inches(0.15), ky + Inches(0.4), Inches(1.9), Inches(0.5),
             size=24, color=kpi_clr, bold=True, font_name="Consolas")


# ============================================================
# SLIDE 8: HERO CHART
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "07", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "HERO CHART SYSTEM", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uD575\uC2EC \uCEF4\uD3EC\uB10C\uD2B8 #2  \u2014  \uC911\uC559 55% \uB300\uD615 \uCC28\uD2B8", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# 3 chart mode cards
modes = [
    ("01  EQUITY CURVE", GREEN,
     ["Portfolio vs KOSPI \uBE44\uAD50", "1D / 1W / 1M \uAE30\uAC04 \uC804\uD658", "\uCD08\uACFC\uC218\uC775\uB960 \uC624\uBC84\uB808\uC774",
      "\uD750\uB9BF\uB0A0 \uC560\uB2C8\uBA54\uC774\uC158"]),
    ("02  PnL WATERFALL", BLUE,
     ["\uC885\uBAA9\uBCC4 \uC218\uC775/\uC190\uC2E4 \uBC14 \uCC28\uD2B8", "PnL% \uAE30\uC900 \uC815\uB82C", "\uD638\uBC84 \uC2DC \uC0C1\uC138\uC815\uBCF4 \uD234\uD301",
      "\uC0C1\uC704/\uD558\uC704 \uC885\uBAA9 \uD558\uC774\uB77C\uC774\uD2B8"]),
    ("03  RISK HEATMAP", RED,
     ["20\uC885\uBAA9 \uACA9\uC790 \uBC30\uCE58", "Trail Gap % \uC0C9\uC0C1 \uCF54\uB529", "Cluster DD \uAC10\uC9C0 \uD45C\uC2DC",
      "Gap <2% \uD3C4\uC2A4 \uC560\uB2C8\uBA54\uC774\uC158"]),
]

for i, (title, accent, bullets) in enumerate(modes):
    cx = Inches(0.6 + i * 4.1)
    cy = Inches(1.4)
    cw = Inches(3.8)
    ch = Inches(3.8)
    add_rounded_rect(slide, cx, cy, cw, ch, BG_CARD, border_color=DIMMER)
    add_rect(slide, cx, cy, cw, Inches(0.05), accent)

    add_text(slide, title, cx + Inches(0.2), cy + Inches(0.2), cw - Inches(0.4), Inches(0.3),
             size=14, color=accent, bold=True, font_name="Consolas")

    # Mock chart area
    add_rounded_rect(slide, cx + Inches(0.15), cy + Inches(0.65), cw - Inches(0.3), Inches(1.5), BG_HOVER)
    if i == 0:
        # Line chart mock
        for j in range(4):
            lx = cx + Inches(0.3 + j * 0.8)
            ly = cy + Inches(1.0 + (j % 2) * 0.4)
            lw = Inches(0.7)
            add_rect(slide, lx, ly, lw, Pt(2), GREEN)
    elif i == 1:
        # Bar chart mock
        heights = [0.8, 0.6, 0.3, -0.2, -0.5, -0.7]
        for j, bh in enumerate(heights):
            bx = cx + Inches(0.35 + j * 0.55)
            if bh > 0:
                by = cy + Inches(1.5 - bh * 0.8)
                add_rect(slide, bx, by, Inches(0.35), Inches(abs(bh) * 0.8), GREEN)
            else:
                by = cy + Inches(1.5)
                add_rect(slide, bx, by, Inches(0.35), Inches(abs(bh) * 0.8), RED)
    else:
        # Heatmap mock
        for r in range(3):
            for c in range(5):
                gx = cx + Inches(0.25 + c * 0.63)
                gy = cy + Inches(0.75 + r * 0.42)
                g_clr = [GREEN, ACCENT_GRN, YELLOW, ORANGE, RED][abs(r * 2 + c) % 5]
                add_rounded_rect(slide, gx, gy, Inches(0.55), Inches(0.35), g_clr)

    # Bullets
    for j, b in enumerate(bullets):
        add_text(slide, f"\u25B8  {b}", cx + Inches(0.2), cy + Inches(2.3 + j * 0.32), cw - Inches(0.4), Inches(0.3),
                 size=10, color=WHITE, font_name="Segoe UI")

# Bottom stats
add_rect(slide, Inches(0.6), Inches(5.5), Inches(12.1), Pt(1), DIMMER)
add_text(slide, "CHART FOOTER:  Top Performer  \u00B7  Bottom Performer  \u00B7  Sharpe Ratio  \u00B7  Win Rate  \u00B7  Excess Return",
         Inches(0.6), Inches(5.7), Inches(12.1), Inches(0.3),
         size=11, color=DIM, font_name="Consolas", align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: ALERT STREAM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "08", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "ALERT STREAM", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uD575\uC2EC \uCEF4\uD3EC\uB10C\uD2B8 #3  \u2014  AI Advisor \uD1B5\uD569 \uC54C\uB9BC", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Left: Priority levels
add_text(slide, "4-LEVEL PRIORITY", Inches(0.6), Inches(1.4), Inches(5), Inches(0.35),
         size=16, color=WHITE, bold=True)

priorities = [
    ("CRITICAL", RED, "\uC989\uC2DC \uD655\uC778 \uD544\uC218", "DD Guard \uBC1C\uB3D9, SAFE_MODE, RECON \uC774\uC0C1"),
    ("HIGH", ORANGE, "\uBE60\uB978 \uD655\uC778 \uD544\uC694", "Trail Stop \uADFC\uC811, Cluster DD, Volume Anomaly"),
    ("MEDIUM", YELLOW, "\uCC38\uACE0 \uC815\uBCF4", "Win Rate \uD558\uB77D, \uD30C\uB77C\uBBF8\uD130 \uCD94\uCC9C, \uC9D1\uC911\uB3C4 \uACBD\uACE0"),
    ("LOW", DIM, "\uBC30\uACBD \uC815\uBCF4", "\uC2DC\uC2A4\uD15C \uC0C1\uD0DC, \uB370\uC774\uD130 \uD488\uC9C8, \uBC30\uCE58 \uACB0\uACFC"),
]

for i, (level, clr, desc, example) in enumerate(priorities):
    py = Inches(1.9 + i * 1.15)
    add_rounded_rect(slide, Inches(0.6), py, Inches(5.5), Inches(0.95), BG_CARD, border_color=DIMMER)
    add_rect(slide, Inches(0.6), py, Inches(0.06), Inches(0.95), clr)

    add_circle(slide, Inches(0.85), py + Inches(0.15), Inches(0.18), clr)
    add_text(slide, level, Inches(1.2), py + Inches(0.1), Inches(2), Inches(0.25),
             size=13, color=clr, bold=True, font_name="Consolas")
    add_text(slide, desc, Inches(3.0), py + Inches(0.12), Inches(3), Inches(0.25),
             size=10, color=WHITE)
    add_text(slide, example, Inches(0.85), py + Inches(0.5), Inches(5.0), Inches(0.3),
             size=9, color=DIM, font_name="Consolas")

# Right: AI Advisor Phases
add_text(slide, "AI ADVISOR 4-PHASE", Inches(6.8), Inches(1.4), Inches(6), Inches(0.35),
         size=16, color=WHITE, bold=True)

phases = [
    ("Phase 1-2", "FACT ANALYSIS", BLUE, ["PnL Attribution", "MDD Contributors", "Entry/Exit Quality", "Operational Issues"]),
    ("Phase 3", "RECOMMENDATIONS", CYAN, ["TRAIL_PCT \uC870\uC815", "REBAL_DAYS \uC81C\uC548", "Strategy Drift \uAC10\uC9C0", "Confidence: LOW/MED"]),
    ("Phase 4", "INTRADAY RISK", RED, ["Flash Drop \uAC10\uC9C0", "Volume Anomaly", "VWAP Divergence", "Cluster DD Alert"]),
]

for i, (phase, title, accent, items) in enumerate(phases):
    px = Inches(6.8)
    py = Inches(1.9 + i * 1.7)
    pw = Inches(5.9)
    ph = Inches(1.5)
    add_rounded_rect(slide, px, py, pw, ph, BG_CARD, border_color=DIMMER)
    add_rect(slide, px, py, pw, Inches(0.04), accent)

    add_text(slide, phase, px + Inches(0.15), py + Inches(0.12), Inches(1.2), Inches(0.25),
             size=10, color=accent, bold=True, font_name="Consolas")
    add_text(slide, title, px + Inches(1.5), py + Inches(0.12), Inches(3), Inches(0.25),
             size=11, color=WHITE, bold=True)

    for j, item in enumerate(items):
        col = j // 2
        row = j % 2
        add_text(slide, f"\u25B8  {item}", px + Inches(0.15 + col * 2.8), py + Inches(0.5 + row * 0.4),
                 Inches(2.7), Inches(0.3), size=9, color=DIM, font_name="Consolas")


# ============================================================
# SLIDE 10: POSITION GRID
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "09", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "POSITION GRID", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uD575\uC2EC \uCEF4\uD3EC\uB10C\uD2B8 #4  \u2014  20\uC885\uBAA9 \uCE74\uB4DC\uD615 \uD0C0\uC77C", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Grid of mock stock tiles (4 rows x 5 cols)
stocks = [
    ("\uC0BC\uC131\uC804\uC790", "+8.2%", "32d", "7.1%", GREEN),
    ("SK\uD558\uC774\uB2C9\uC2A4", "+5.7%", "28d", "5.3%", GREEN),
    ("\uD604\uB300\uCC28", "+4.1%", "15d", "8.2%", GREEN),
    ("LG\uC5D0\uB108\uC9C0", "+3.3%", "22d", "4.8%", GREEN),
    ("POSCO", "+2.8%", "31d", "6.4%", GREEN),
    ("\uCE74\uCE74\uC624", "+2.1%", "18d", "3.2%", ACCENT_GRN),
    ("\uB124\uC774\uBC84", "+1.5%", "12d", "9.1%", ACCENT_GRN),
    ("KB\uAE08\uC735", "+0.8%", "25d", "11.2%", WHITE),
    ("\uC0BC\uC131\uBB3C\uC0B0", "+0.3%", "9d", "10.5%", WHITE),
    ("\uD55C\uD654\uC194\uB8E8\uC158", "-0.2%", "7d", "8.7%", WHITE),
    ("LG\uD654\uD559", "-0.9%", "14d", "5.5%", ORANGE),
    ("\uAE30\uC544", "-1.5%", "20d", "3.8%", ORANGE),
    ("\uD604\uB300\uBAA8\uBE44\uC2A4", "-2.1%", "11d", "2.9%", RED),
    ("S-Oil", "-2.8%", "26d", "2.1%", RED),
    ("\uB450\uC0B0", "-3.5%", "16d", "1.4%", RED),
    ("\uD55C\uC804KPS", "+1.2%", "19d", "6.8%", ACCENT_GRN),
    ("\uC0BC\uC131SDI", "+4.5%", "24d", "5.0%", GREEN),
    ("\uCCAD\uD638\uB098\uC774\uC2A4", "-0.5%", "8d", "9.3%", WHITE),
    ("KT&G", "+1.8%", "30d", "7.7%", ACCENT_GRN),
    ("\uD55C\uAD6D\uC804\uB825", "-1.2%", "13d", "4.1%", ORANGE),
]

for i, (name, pnl, hold, gap, clr) in enumerate(stocks):
    row = i // 5
    col = i % 5
    tx = Inches(0.6 + col * 2.45)
    ty = Inches(1.4 + row * 1.35)
    tw = Inches(2.3)
    th = Inches(1.2)

    # Card bg with tint
    if clr == GREEN or clr == ACCENT_GRN:
        card_bg = RGBColor(0x08, 0x1A, 0x12)
    elif clr == RED:
        card_bg = RGBColor(0x1A, 0x08, 0x0A)
    elif clr == ORANGE:
        card_bg = RGBColor(0x1A, 0x12, 0x08)
    else:
        card_bg = BG_CARD

    add_rounded_rect(slide, tx, ty, tw, th, card_bg, border_color=DIMMER)
    # Top accent
    add_rect(slide, tx, ty, tw, Inches(0.03), clr)

    # Stock name
    add_text(slide, name, tx + Inches(0.1), ty + Inches(0.08), tw - Inches(0.2), Inches(0.22),
             size=10, color=WHITE, bold=True, font_name="Segoe UI")
    # PnL
    add_text(slide, pnl, tx + Inches(0.1), ty + Inches(0.35), Inches(1.0), Inches(0.35),
             size=18, color=clr, bold=True, font_name="Consolas")
    # Hold days & Gap
    add_text(slide, f"Hold {hold}", tx + Inches(1.2), ty + Inches(0.38), Inches(0.8), Inches(0.2),
             size=8, color=DIM, font_name="Consolas")
    add_text(slide, f"Gap {gap}", tx + Inches(1.2), ty + Inches(0.58), Inches(0.8), Inches(0.2),
             size=8, color=YELLOW if float(gap.replace('%','')) < 3 else DIM, font_name="Consolas")

    # Warning indicator for near trail
    if float(gap.replace('%','')) < 3:
        add_circle(slide, tx + tw - Inches(0.3), ty + Inches(0.08), Inches(0.15), RED)
        add_text(slide, "!", tx + tw - Inches(0.3), ty + Inches(0.06), Inches(0.15), Inches(0.15),
                 size=9, color=WHITE, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)

# Sort options
add_text(slide, "SORT:  PnL%  \u00B7  Trail Gap  \u00B7  Hold Days  \u00B7  Sector",
         Inches(0.6), Inches(6.9), Inches(12), Inches(0.3),
         size=11, color=DIM, font_name="Consolas")


# ============================================================
# SLIDE 11: SYSTEM HEARTBEAT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "10", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "SYSTEM HEARTBEAT", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uD575\uC2EC \uCEF4\uD3EC\uB10C\uD2B8 #5  \u2014  \uC2DC\uC2A4\uD15C \uC0DD\uBA85\uC120", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Normal mode (minimal 8px)
add_text(slide, "NORMAL MODE  \u2014  Minimal (8px)", Inches(0.6), Inches(1.4), Inches(6), Inches(0.3),
         size=14, color=GREEN, bold=True, font_name="Consolas")
add_rounded_rect(slide, Inches(0.6), Inches(1.8), Inches(12.1), Inches(0.4), BG_CARD, border_color=GREEN)
for i, (lbl, clr) in enumerate([("ENGINE", GREEN), ("STALE", GREEN), ("RISK", GREEN), ("REGIME", BLUE), ("15:32:01", DIM)]):
    lx = Inches(0.9 + i * 2.4)
    add_circle(slide, lx, Inches(1.88), Inches(0.14), clr)
    add_text(slide, lbl, lx + Inches(0.22), Inches(1.85), Inches(1.5), Inches(0.22),
             size=9, color=DIM, font_name="Consolas", bold=True)

# Alert mode (expanded 40px)
add_text(slide, "ALERT MODE  \u2014  Expanded (40px)", Inches(0.6), Inches(2.6), Inches(6), Inches(0.3),
         size=14, color=RED, bold=True, font_name="Consolas")
add_rounded_rect(slide, Inches(0.6), Inches(3.0), Inches(12.1), Inches(1.0), RGBColor(0x1A, 0x08, 0x0A), border_color=RED)
for i, (lbl, clr) in enumerate([("ENGINE", GREEN), ("STALE", RED), ("RISK", RED), ("REGIME", YELLOW), ("15:32:01", RED)]):
    lx = Inches(0.9 + i * 2.4)
    add_circle(slide, lx, Inches(3.15), Inches(0.14), clr)
    add_text(slide, lbl, lx + Inches(0.22), Inches(3.12), Inches(1.5), Inches(0.22),
             size=9, color=clr, font_name="Consolas", bold=True)
add_text(slide, "\u26A0  STALE DATA (243s)  |  RISK: DD_SEVERE  |  BUY BLOCKED",
         Inches(0.9), Inches(3.5), Inches(11.5), Inches(0.3),
         size=12, color=RED, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)

# Risk mode ladder
add_text(slide, "RISK MODE \uACC4\uCE35 (8\uB2E8\uACC4)", Inches(0.6), Inches(4.4), Inches(6), Inches(0.3),
         size=14, color=WHITE, bold=True)

risk_levels = [
    ("NORMAL", GREEN), ("DD_CAUTION", YELLOW), ("DD_WARNING", YELLOW),
    ("DD_SEVERE", ORANGE), ("DD_CRITICAL", RED),
    ("SAFE_L1", RED), ("SAFE_L2", RED), ("BLOCKED", RED),
]
for i, (rl, rc) in enumerate(risk_levels):
    rx = Inches(0.6 + i * 1.5)
    ry = Inches(4.8)
    add_rounded_rect(slide, rx, ry, Inches(1.35), Inches(0.5), BG_CARD, border_color=rc)
    add_text(slide, rl, rx + Inches(0.05), ry + Inches(0.08), Inches(1.25), Inches(0.35),
             size=8, color=rc, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)
    if i < 7:
        add_text(slide, "\u2192", rx + Inches(1.35), ry + Inches(0.08), Inches(0.15), Inches(0.35),
                 size=10, color=DIMMER, font_name="Consolas", align=PP_ALIGN.CENTER)

# Stale detection
add_text(slide, "STALE DETECTION", Inches(0.6), Inches(5.6), Inches(6), Inches(0.3),
         size=14, color=WHITE, bold=True)
stale_data = [
    ("<90s", "FRESH", GREEN, "\uC815\uC0C1 \uC6B4\uC601"),
    ("90~180s", "WARN", YELLOW, "\uC8FC\uC758 \u2014 BUY LIMITED"),
    (">180s", "STALE", RED, "\uC704\uD5D8 \u2014 BUY BLOCKED + \uBC30\uB108"),
]
for i, (time, status, clr, desc) in enumerate(stale_data):
    sx = Inches(0.6 + i * 4.1)
    sy = Inches(6.0)
    add_rounded_rect(slide, sx, sy, Inches(3.8), Inches(0.7), BG_CARD, border_color=clr)
    add_rect(slide, sx, sy, Inches(0.05), Inches(0.7), clr)
    add_text(slide, f"{time}  {status}", sx + Inches(0.2), sy + Inches(0.08), Inches(3.4), Inches(0.25),
             size=12, color=clr, bold=True, font_name="Consolas")
    add_text(slide, desc, sx + Inches(0.2), sy + Inches(0.38), Inches(3.4), Inches(0.25),
             size=10, color=DIM)


# ============================================================
# SLIDE 12: COLOR SYSTEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "11", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "COLOR SYSTEM & TYPOGRAPHY", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uCEEC\uB7EC \uC2DC\uC2A4\uD15C & \uD0C0\uC774\uD3EC\uADF8\uB798\uD53C", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Signal colors
add_text(slide, "SIGNAL PALETTE", Inches(0.6), Inches(1.4), Inches(5), Inches(0.3),
         size=14, color=WHITE, bold=True)

signal_colors = [
    ("Signal Green", "#00FF88", GREEN, "\uC218\uC775, \uC815\uC0C1, BUY OK"),
    ("Alert Red", "#FF3344", RED, "\uC190\uC2E4, \uC704\uD5D8, BLOCKED"),
    ("Caution Yellow", "#FFDD00", YELLOW, "\uC8FC\uC758, LIMITED"),
    ("Info Blue", "#6699FF", BLUE, "\uC911\uB9BD, KOSPI, \uCC38\uACE0"),
    ("Cyan", "#44CCCC", CYAN, "\uBCF4\uC870"),
]

for i, (name, hex_val, clr, usage) in enumerate(signal_colors):
    cx = Inches(0.6 + i * 2.45)
    cy = Inches(1.8)
    # Color swatch
    add_rounded_rect(slide, cx, cy, Inches(2.2), Inches(0.8), clr)
    add_text(slide, hex_val, cx + Inches(0.1), cy + Inches(0.1), Inches(2.0), Inches(0.35),
             size=14, color=BG_MAIN, bold=True, font_name="Consolas")
    add_text(slide, name, cx, cy + Inches(0.85), Inches(2.2), Inches(0.25),
             size=10, color=WHITE, bold=True)
    add_text(slide, usage, cx, cy + Inches(1.1), Inches(2.2), Inches(0.25),
             size=9, color=DIM)

# Background layers
add_text(slide, "BACKGROUND LAYERS", Inches(0.6), Inches(3.3), Inches(5), Inches(0.3),
         size=14, color=WHITE, bold=True)

bg_layers = [
    ("L0  Main", "#060810", BG_MAIN),
    ("L1  Card", "#0E1118", BG_CARD),
    ("L2  Hover", "#161A24", BG_HOVER),
    ("L3  Active", "#1E2230", BG_ACTIVE),
]

for i, (name, hex_val, clr) in enumerate(bg_layers):
    bx = Inches(0.6 + i * 3.05)
    by = Inches(3.7)
    add_rounded_rect(slide, bx, by, Inches(2.8), Inches(1.0), clr, border_color=DIM)
    add_text(slide, hex_val, bx + Inches(0.15), by + Inches(0.15), Inches(2.5), Inches(0.35),
             size=14, color=WHITE, bold=True, font_name="Consolas")
    add_text(slide, name, bx + Inches(0.15), by + Inches(0.55), Inches(2.5), Inches(0.3),
             size=11, color=DIM)

# Typography
add_text(slide, "TYPOGRAPHY HIERARCHY", Inches(0.6), Inches(5.0), Inches(5), Inches(0.3),
         size=14, color=WHITE, bold=True)

typo_data = [
    ("NUMBERS & DATA", "Consolas  \u00B7  Monospace", Inches(2.8)),
    ("LABELS & UI", "Segoe UI  \u00B7  Sans-serif", Inches(2.8)),
    ("HEADINGS", "Segoe UI Light  \u00B7  28~40pt", Inches(2.8)),
]

for i, (role, font_desc, tw) in enumerate(typo_data):
    tx = Inches(0.6 + i * 4.1)
    ty = Inches(5.4)
    add_rounded_rect(slide, tx, ty, Inches(3.8), Inches(1.3), BG_CARD, border_color=DIMMER)
    add_text(slide, role, tx + Inches(0.15), ty + Inches(0.1), Inches(3.5), Inches(0.25),
             size=10, color=GREEN, bold=True, font_name="Consolas")
    add_text(slide, font_desc, tx + Inches(0.15), ty + Inches(0.4), Inches(3.5), Inches(0.25),
             size=10, color=WHITE)
    # Sample
    if i == 0:
        add_text(slide, "+47.3%", tx + Inches(0.15), ty + Inches(0.7), Inches(3.5), Inches(0.4),
                 size=24, color=GREEN, bold=True, font_name="Consolas")
    elif i == 1:
        add_text(slide, "Today PnL  \u00B7  Current MDD", tx + Inches(0.15), ty + Inches(0.7), Inches(3.5), Inches(0.4),
                 size=14, color=WHITE, font_name="Segoe UI")
    else:
        add_text(slide, "Q-TRON CONSOLE", tx + Inches(0.15), ty + Inches(0.7), Inches(3.5), Inches(0.4),
                 size=20, color=WHITE, font_name="Segoe UI Light")


# ============================================================
# SLIDE 13: INTERACTION DESIGN
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "12", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "INTERACTION DESIGN", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uC778\uD130\uB799\uC158 \uC124\uACC4  \u2014  READ-ONLY \uC6D0\uCE59", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Safety banner
add_rounded_rect(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(0.6), RGBColor(0x33, 0x0A, 0x0E), border_color=RED)
add_text(slide, "\u26A0  READ-ONLY PRINCIPLE  \u2014  \uBAA8\uB2C8\uD130\uB294 \uC808\uB300 \uC8FC\uBB38\uC744 \uC2E4\uD589\uD560 \uC218 \uC5C6\uC74C  |  \uAD00\uCC30 \uC804\uC6A9, \uC81C\uC5B4 \uBD88\uAC00",
         Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.4),
         size=13, color=RED, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)

# 4 interaction types
interactions = [
    ("HOVER", "\uD234\uD301 \uC2DC\uC2A4\uD15C", CYAN,
     ["\uCC28\uD2B8 \uD638\uBC84: \uC885\uBAA9 \uC0C1\uC138 \uD234\uD301", "\uD14C\uC774\uBE14 \uD638\uBC84: \uD589 \uD558\uC774\uB77C\uC774\uD2B8",
      "\uCE74\uB4DC \uD638\uBC84: \uC804\uCCB4 \uC815\uBCF4 \uD45C\uC2DC", "Trail Gap: \uC9C4\uC785\uAC00/\uD604\uC7AC\uAC00/\uC13C\uD130"]),
    ("CLICK", "\uBAA8\uB4DC \uC804\uD658", BLUE,
     ["\uCC28\uD2B8 \uD0ED: Equity/PnL/Heatmap", "\uC815\uB82C: PnL%/Gap/Hold Days",
      "\uC54C\uB9BC: \uC0C1\uC138 \uD3BC\uCE68/\uC811\uAE30", "KPI: \uC0C1\uC138 \uD31D\uC5C5 \uD1A0\uAE00"]),
    ("KEYBOARD", "\uB2E8\uCD95\uD0A4", GREEN,
     ["1/2/3: \uCC28\uD2B8 \uBAA8\uB4DC \uC804\uD658", "R: \uC989\uC2DC \uC0C8\uB85C\uACE0\uCE68",
      "F: \uC804\uCCB4\uD654\uBA74 \uD1A0\uAE00", "A: Advisor \uD328\uB110 \uD1A0\uAE00"]),
    ("ANIMATION", "\uC54C\uB9BC \uD750\uB984", YELLOW,
     ["\uC2E0\uADDC \u2192 \uC2AC\uB77C\uC774\uB4DC\uC778 \uB4F1\uC7A5", "3\uCD08 \uD558\uC774\uB77C\uC774\uD2B8 \uAE5C\uBC15",
      "\uC2A4\uD2B8\uB9BC \uD569\uB958 \u2192 \uC815\uB82C", "Trail <2% \uD3C4\uC2A4 \uC560\uB2C8\uBA54\uC774\uC158"]),
]

for i, (title, ko, accent, items) in enumerate(interactions):
    cx = Inches(0.6 + i * 3.15)
    cy = Inches(2.3)
    cw = Inches(2.95)
    ch = Inches(4.5)
    add_rounded_rect(slide, cx, cy, cw, ch, BG_CARD, border_color=DIMMER)
    add_rect(slide, cx, cy, cw, Inches(0.04), accent)

    add_text(slide, title, cx + Inches(0.15), cy + Inches(0.2), cw - Inches(0.3), Inches(0.3),
             size=16, color=accent, bold=True, font_name="Consolas")
    add_text(slide, ko, cx + Inches(0.15), cy + Inches(0.5), cw - Inches(0.3), Inches(0.25),
             size=11, color=WHITE)

    add_rect(slide, cx + Inches(0.15), cy + Inches(0.85), cw - Inches(0.3), Pt(1), DIMMER)

    for j, item in enumerate(items):
        add_text(slide, f"\u25B8  {item}", cx + Inches(0.15), cy + Inches(1.1 + j * 0.7), cw - Inches(0.3), Inches(0.6),
                 size=10, color=DIM, font_name="Segoe UI", line_spacing=14)


# ============================================================
# SLIDE 14: RESPONSIVE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "13", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "RESPONSIVE & ADAPTIVE", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uBC18\uC751\uD615 & \uC801\uC751\uD615 \uB808\uC774\uC544\uC6C3", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# 3 layout modes
modes_resp = [
    ("FULL", "1920 x 1080+", GREEN, "\uC804\uCCB4 Mission Control",
     ["Decision Hub + KPI Stack", "Hero Chart (3\uBAA8\uB4DC)", "Alert Stream + Advisor",
      "Position Grid 20\uC885\uBAA9", "System Heartbeat"]),
    ("COMPACT", "1280 x 720", YELLOW, "2-Column \uC555\uCD95",
     ["Decision + Chart (L)", "Alerts + Positions (R)", "KPI \uC218\uD3C9 \uBC30\uCE58",
      "\uCC28\uD2B8 1\uBAA8\uB4DC\uB9CC", "\uC54C\uB9BC \uCD95\uC18C"]),
    ("MINI", "800 x 600", CYAN, "PIP \uBAA8\uB4DC",
     ["Decision Card \uB2E8\uB3C5", "KPI 5\uAC1C\uB9CC", "\uCD5C\uC18C \uC815\uBCF4",
      "\uD56D\uC0C1 \uCD5C\uC0C1\uC704", "\uBCF4\uC870 \uBAA8\uB2C8\uD130\uC6A9"]),
]

for i, (name, res, accent, desc, items) in enumerate(modes_resp):
    cx = Inches(0.6 + i * 4.1)
    cy = Inches(1.4)
    cw = Inches(3.8)
    ch = Inches(4.0)
    add_rounded_rect(slide, cx, cy, cw, ch, BG_CARD, border_color=DIMMER)
    add_rect(slide, cx, cy, cw, Inches(0.05), accent)

    # Resolution badge
    add_rounded_rect(slide, cx + Inches(0.15), cy + Inches(0.2), Inches(1.8), Inches(0.35), BG_HOVER, border_color=accent)
    add_text(slide, res, cx + Inches(0.15), cy + Inches(0.22), Inches(1.8), Inches(0.3),
             size=11, color=accent, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)

    add_text(slide, name, cx + Inches(2.1), cy + Inches(0.2), Inches(1.5), Inches(0.3),
             size=18, color=WHITE, bold=True, font_name="Consolas")
    add_text(slide, desc, cx + Inches(0.15), cy + Inches(0.7), cw - Inches(0.3), Inches(0.3),
             size=12, color=WHITE)

    for j, item in enumerate(items):
        add_text(slide, f"\u25B8  {item}", cx + Inches(0.15), cy + Inches(1.2 + j * 0.5),
                 cw - Inches(0.3), Inches(0.4), size=10, color=DIM, font_name="Segoe UI")

# Additional features
add_rect(slide, Inches(0.6), Inches(5.7), Inches(12.1), Pt(1), DIMMER)
features = [
    ("\uB2E4\uC911 \uBAA8\uB2C8\uD130", "\uCC28\uD2B8 \uBD84\uB9AC \uAC00\uB2A5 (QDockWidget Float)"),
    ("\uC57C\uAC04 \uBAA8\uB4DC", "\uB354 \uC5B4\uB450\uC6B4 \uBC30\uACBD + \uBC1D\uAE30 \uAC10\uC18C"),
    ("Dock \uC7AC\uBC30\uCE58", "\uD30C\uC6CC\uC720\uC800 \uCEE4\uC2A4\uD130\uB9C8\uC774\uC9D5 \uC720\uC9C0"),
]
for i, (feat, desc) in enumerate(features):
    fx = Inches(0.6 + i * 4.1)
    add_text(slide, feat, fx, Inches(5.9), Inches(3.8), Inches(0.3),
             size=12, color=GREEN, bold=True, font_name="Consolas")
    add_text(slide, desc, fx, Inches(6.2), Inches(3.8), Inches(0.3),
             size=10, color=DIM)


# ============================================================
# SLIDE 15: AI ADVISOR INTEGRATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "14", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "AI ADVISOR INTEGRATION", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "AI \uC5B4\uB4DC\uBC14\uC774\uC800 \uD1B5\uD569 \uC124\uACC4", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Pipeline diagram
phases_adv = [
    ("Phase 1", "INGESTION", BLUE, "Snapshot\nBuilder"),
    ("Phase 2", "ANALYSIS", CYAN, "PnL / MDD\nEntry / Exit"),
    ("Phase 3", "RECOMMEND", GREEN, "Param Tuning\nDrift Guard"),
    ("Phase 4", "INTRADAY", ORANGE, "Flash Drop\nVWAP / Vol"),
    ("Phase 5", "LLM (Future)", DIM, "\uC790\uC5F0\uC5B4\n\uC694\uC57D"),
]

# Pipeline flow
for i, (phase, name, accent, desc) in enumerate(phases_adv):
    px = Inches(0.4 + i * 2.55)
    py = Inches(1.4)
    pw = Inches(2.3)
    ph = Inches(1.5)
    add_rounded_rect(slide, px, py, pw, ph, BG_CARD, border_color=accent)
    add_rect(slide, px, py, pw, Inches(0.04), accent)

    add_text(slide, phase, px + Inches(0.1), py + Inches(0.12), pw - Inches(0.2), Inches(0.2),
             size=9, color=accent, bold=True, font_name="Consolas")
    add_text(slide, name, px + Inches(0.1), py + Inches(0.35), pw - Inches(0.2), Inches(0.3),
             size=13, color=WHITE, bold=True, font_name="Consolas")
    add_text(slide, desc, px + Inches(0.1), py + Inches(0.75), pw - Inches(0.2), Inches(0.6),
             size=10, color=DIM, line_spacing=15)

    if i < 4:
        add_text(slide, "\u25B6", px + pw + Inches(0.02), py + Inches(0.5), Inches(0.2), Inches(0.3),
                 size=14, color=DIMMER, font_name="Consolas")

# GUI integration details
add_text(slide, "GUI \uD1B5\uD569 \uBC29\uBC95", Inches(0.6), Inches(3.2), Inches(5), Inches(0.3),
         size=16, color=WHITE, bold=True)

gui_features = [
    ("SLIDE-OUT PANEL", "\uC6B0\uCE21 \uD655\uC7A5\uD615 Advisor \uC804\uC6A9 \uD328\uB110",
     CYAN, "\uD3C9\uC0C1\uC2DC \uC228\uACA8\uC838 \uC788\uB2E4\uAC00 \uD074\uB9AD/A\uD0A4\uB85C \uD65C\uC131\uD654\n\uC804\uCCB4 Alert + Recommendation \uBAA9\uB85D \uD45C\uC2DC\n\uD544\uD130: \uB0A0\uC9DC, \uC6B0\uC120\uC21C\uC704, \uCE74\uD14C\uACE0\uB9AC"),
    ("ALERT CARDS", "\uCE74\uB4DC \uC2A4\uD0DD (Max 8, \uC6B0\uC120\uC21C\uC704 \uC815\uB82C)",
     ORANGE, "Phase 2 \uACB0\uACFC\uBB3C \uC0C1\uC704 8\uAC74 \uD45C\uC2DC\n\uAE34\uAE09\uB3C4 \uC0C9\uC0C1 \uCF54\uB529 (RED/ORANGE/YELLOW/GRAY)\n\uD074\uB9AD \uC2DC debug_hint \uD45C\uC2DC"),
    ("RECOMMENDATION", "\"\uC81C\uC548\" \uBC30\uC9C0 + \uC0C1\uC138 \uD3BC\uCE68",
     GREEN, "Phase 3 \uD30C\uB77C\uBBF8\uD130 \uCD94\uCC9C \uCE74\uB4DC\nConfidence: LOW/MEDIUM \uD45C\uC2DC\nDrift Guard: \uC2A4\uD338 \uBC29\uC9C0 \uD544\uD130\uB9C1"),
    ("INTRADAY TIMELINE", "\uC2E4\uC2DC\uAC04 \uC704\uD5D8 \uD0C0\uC784\uB77C\uC778",
     RED, "Phase 4 \uC7A5\uC911 \uC774\uBCA4\uD2B8 \uC2DC\uAC04\uC21C \uBC30\uCE58\nFlash Drop, Volume Spike, Near Trail\nRisk Score \uC885\uD569 \uAC8C\uC774\uC9C0"),
]

for i, (title, sub, accent, desc) in enumerate(gui_features):
    gx = Inches(0.6 + (i % 2) * 6.3)
    gy = Inches(3.7 + (i // 2) * 1.85)
    gw = Inches(6.0)
    gh = Inches(1.65)
    add_rounded_rect(slide, gx, gy, gw, gh, BG_CARD, border_color=DIMMER)
    add_rect(slide, gx, gy, Inches(0.05), gh, accent)

    add_text(slide, title, gx + Inches(0.2), gy + Inches(0.1), gw - Inches(0.4), Inches(0.25),
             size=12, color=accent, bold=True, font_name="Consolas")
    add_text(slide, sub, gx + Inches(0.2), gy + Inches(0.35), gw - Inches(0.4), Inches(0.25),
             size=10, color=WHITE)
    add_text(slide, desc, gx + Inches(0.2), gy + Inches(0.65), gw - Inches(0.4), Inches(0.9),
             size=9, color=DIM, line_spacing=14)


# ============================================================
# SLIDE 16: RISK VISUALIZATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "15", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "RISK VISUALIZATION", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uB9AC\uC2A4\uD06C \uC2DC\uAC01\uD654 \uD601\uC2E0", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# 4 risk viz concepts
risk_viz = [
    ("DD GUARD GAUGE", "\uAC8C\uC774\uC9C0 \uCC28\uD2B8", RED,
     "\uC77C -4% / \uC6D4 -7% \uC784\uACC4\uC810 \uD45C\uC2DC\n\uD604\uC7AC DD \uC218\uC900 \uC2E4\uC2DC\uAC04 \uBC18\uC601\n\uC784\uACC4 \uC811\uADFC \uC2DC \uC560\uB2C8\uBA54\uC774\uC158 \uACBD\uACE0"),
    ("TRAIL STOP RADAR", "\uB808\uC774\uB354 \uCC28\uD2B8", YELLOW,
     "20\uC885\uBAA9 \uC6D0\uD615 \uBC30\uCE58 (\uD504\uB85C\uD3A0\uB9AC\uC624 \uB808\uC774\uB354)\n\uC911\uC2EC \uAC00\uAE4C\uC6B8\uC218\uB85D \uC704\uD5D8 (Gap %)\n\uBE68\uAC04 \uC601\uC5ED = Stop Loss \uC784\uBC15"),
    ("RISK TIMELINE", "\uC2DC\uAC04\uBCC4 \uB9AC\uC2A4\uD06C \uBCC0\uD654", BLUE,
     "\uD558\uB8E8 \uC911 \uB9AC\uC2A4\uD06C \uB808\uBCA8 \uBCC0\uD654 \uD0C0\uC784\uB77C\uC778\n09:00~15:30 \uAD6C\uAC04\uBCC4 \uC0C1\uD0DC \uD45C\uC2DC\nSAFE_MODE \uC804\uD658 \uC2DC\uC810 \uB9C8\uCEE4"),
    ("CLUSTER DD ALERT", "\uD074\uB7EC\uC2A4\uD130 DD \uAC10\uC9C0", ORANGE,
     "3+ \uC885\uBAA9 \uB3D9\uC2DC -3% \u2192 \uBE68\uAC04 \uC624\uBC84\uB808\uC774\nSAFE_MODE \uC804\uD658 \uC804\uCCB4 \uD654\uBA74 \uD1A4 \uBCC0\uACBD\n\uC815\uC0C1=\uBE14\uB8E8\uD2F4\uD2B8, \uC704\uD5D8=\uB808\uB4DC\uD2F4\uD2B8"),
]

for i, (title, sub, accent, desc) in enumerate(risk_viz):
    cx = Inches(0.6 + (i % 2) * 6.3)
    cy = Inches(1.4 + (i // 2) * 2.8)
    cw = Inches(6.0)
    ch = Inches(2.5)
    add_rounded_rect(slide, cx, cy, cw, ch, BG_CARD, border_color=DIMMER)
    add_rect(slide, cx, cy, cw, Inches(0.05), accent)

    # Mock viz area
    viz_x = cx + Inches(0.15)
    viz_y = cy + Inches(0.15)
    viz_w = Inches(2.0)
    viz_h = Inches(2.2)
    add_rounded_rect(slide, viz_x, viz_y, viz_w, viz_h, BG_HOVER)

    if i == 0:  # Gauge
        # Semi-circle gauge mock
        add_circle(slide, viz_x + Inches(0.35), viz_y + Inches(0.5), Inches(1.2), BG_ACTIVE)
        add_circle(slide, viz_x + Inches(0.55), viz_y + Inches(0.7), Inches(0.8), BG_HOVER)
        add_text(slide, "-2.1%", viz_x + Inches(0.35), viz_y + Inches(0.8), Inches(1.2), Inches(0.4),
                 size=16, color=YELLOW, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)
        add_text(slide, "Daily DD", viz_x + Inches(0.35), viz_y + Inches(1.2), Inches(1.2), Inches(0.3),
                 size=9, color=DIM, font_name="Consolas", align=PP_ALIGN.CENTER)
    elif i == 1:  # Radar
        add_circle(slide, viz_x + Inches(0.4), viz_y + Inches(0.4), Inches(1.1), BG_ACTIVE)
        add_circle(slide, viz_x + Inches(0.6), viz_y + Inches(0.6), Inches(0.7), BG_HOVER)
        add_circle(slide, viz_x + Inches(0.8), viz_y + Inches(0.8), Inches(0.3), RGBColor(0x33, 0x0A, 0x0E))
        # Stock dots
        import random
        random.seed(42)
        for _ in range(8):
            dx = viz_x + Inches(random.uniform(0.3, 1.6))
            dy = viz_y + Inches(random.uniform(0.3, 1.6))
            add_circle(slide, dx, dy, Inches(0.08), GREEN)
        for _ in range(3):
            dx = viz_x + Inches(random.uniform(0.7, 1.2))
            dy = viz_y + Inches(random.uniform(0.7, 1.2))
            add_circle(slide, dx, dy, Inches(0.08), RED)
    elif i == 2:  # Timeline
        for j in range(7):
            tx = viz_x + Inches(0.15 + j * 0.25)
            th = [0.8, 0.4, 0.3, 0.5, 1.0, 0.6, 0.3][j]
            tc = [GREEN, GREEN, GREEN, YELLOW, RED, YELLOW, GREEN][j]
            add_rect(slide, tx, viz_y + Inches(1.8 - th), Inches(0.18), Inches(th), tc)
        add_text(slide, "09:00        15:30", viz_x + Inches(0.1), viz_y + Inches(1.85), Inches(1.8), Inches(0.2),
                 size=7, color=DIM, font_name="Consolas")
    else:  # Cluster
        for r in range(4):
            for c in range(4):
                gx = viz_x + Inches(0.1 + c * 0.45)
                gy = viz_y + Inches(0.1 + r * 0.5)
                is_cluster = (r >= 2 and c >= 2)
                g_clr = RED if is_cluster else GREEN
                add_rounded_rect(slide, gx, gy, Inches(0.38), Inches(0.4), g_clr)

    # Text content
    add_text(slide, title, cx + Inches(2.3), cy + Inches(0.2), Inches(3.5), Inches(0.3),
             size=14, color=accent, bold=True, font_name="Consolas")
    add_text(slide, sub, cx + Inches(2.3), cy + Inches(0.5), Inches(3.5), Inches(0.25),
             size=11, color=WHITE)
    add_text(slide, desc, cx + Inches(2.3), cy + Inches(0.9), Inches(3.5), Inches(1.2),
             size=10, color=DIM, line_spacing=16)


# ============================================================
# SLIDE 17: TECHNICAL FEASIBILITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

add_rect(slide, Inches(0), Inches(0), SW, Inches(1.1), BG_CARD)
add_text(slide, "16", Inches(0.6), Inches(0.2), Inches(1), Inches(0.6),
         size=32, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "TECHNICAL FEASIBILITY", Inches(1.5), Inches(0.15), Inches(8), Inches(0.5),
         size=28, color=WHITE, bold=True, font_name="Segoe UI")
add_text(slide, "\uAE30\uC220 \uAD6C\uD604 \uAC00\uB2A5\uC131", Inches(1.5), Inches(0.6), Inches(8), Inches(0.4),
         size=14, color=DIM, font_name="Segoe UI")

# Stack comparison
stacks = [
    ("CURRENT", "PyQt5 + pyqtgraph", DIM,
     ["QDockWidget \uAE30\uBC18", "pyqtgraph \uCC28\uD2B8", "QSS \uC2A4\uD0C0\uC77C\uB9C1", "Python 3.9 (32-bit)"],
     "\uD604\uC7AC \uC6B4\uC601 \uC911"),
    ("PROPOSED", "PyQt6 + pyqtgraph", GREEN,
     ["QDockWidget \uC720\uC9C0", "pyqtgraph + Custom Widget", "QSS \uACE0\uAE09 \uD14C\uB9C8", "Python 3.13 (64-bit)"],
     "\uCD94\uCC9C \u2014 \uCD5C\uC18C \uBCC0\uACBD"),
    ("ALTERNATIVE", "Electron + React + D3", BLUE,
     ["Web \uAE30\uBC18 \uD655\uC7A5\uC131", "D3.js \uACE0\uAE09 \uCC28\uD2B8", "CSS \uC560\uB2C8\uBA54\uC774\uC158", "\uB9AC\uBAA8\uD2B8 \uC811\uC18D \uAC00\uB2A5"],
     "\uC7A5\uAE30 \uBE44\uC804 (Gen5)"),
]

for i, (name, stack, accent, items, note) in enumerate(stacks):
    cx = Inches(0.6 + i * 4.1)
    cy = Inches(1.4)
    cw = Inches(3.8)
    ch = Inches(3.2)
    add_rounded_rect(slide, cx, cy, cw, ch, BG_CARD, border_color=DIMMER)
    add_rect(slide, cx, cy, cw, Inches(0.05), accent)

    add_text(slide, name, cx + Inches(0.2), cy + Inches(0.15), cw - Inches(0.4), Inches(0.3),
             size=16, color=accent, bold=True, font_name="Consolas")
    add_text(slide, stack, cx + Inches(0.2), cy + Inches(0.45), cw - Inches(0.4), Inches(0.25),
             size=11, color=WHITE, bold=True)

    for j, item in enumerate(items):
        add_text(slide, f"\u25B8  {item}", cx + Inches(0.2), cy + Inches(0.9 + j * 0.4),
                 cw - Inches(0.4), Inches(0.35), size=10, color=DIM, font_name="Segoe UI")

    add_rounded_rect(slide, cx + Inches(0.2), cy + ch - Inches(0.45), cw - Inches(0.4), Inches(0.3), BG_HOVER, border_color=accent)
    add_text(slide, note, cx + Inches(0.2), cy + ch - Inches(0.43), cw - Inches(0.4), Inches(0.25),
             size=9, color=accent, font_name="Consolas", align=PP_ALIGN.CENTER, bold=True)

# Performance targets
add_rect(slide, Inches(0.6), Inches(4.8), Inches(12.1), Pt(1), DIMMER)
add_text(slide, "PERFORMANCE TARGETS", Inches(0.6), Inches(5.0), Inches(5), Inches(0.3),
         size=14, color=WHITE, bold=True)

targets = [
    ("60fps", "\uCC28\uD2B8 \uB80C\uB354\uB9C1"),
    ("<100ms", "\uB370\uC774\uD130 \uAC31\uC2E0"),
    ("<50MB", "RAM \uC0AC\uC6A9\uB7C9"),
    ("3s", "\uD3F4\uB9C1 \uC8FC\uAE30"),
]
for i, (val, desc) in enumerate(targets):
    tx = Inches(0.6 + i * 3.05)
    add_text(slide, val, tx, Inches(5.4), Inches(2.8), Inches(0.5),
             size=32, color=GREEN, bold=True, font_name="Consolas")
    add_text(slide, desc, tx, Inches(5.9), Inches(2.8), Inches(0.3),
             size=11, color=DIM)

# Implementation timeline
add_text(slide, "IMPLEMENTATION TIMELINE", Inches(0.6), Inches(6.3), Inches(5), Inches(0.3),
         size=14, color=WHITE, bold=True)

timeline = [
    ("Phase 1", "4\uC8FC", "\uD575\uC2EC \uB808\uC774\uC544\uC6C3 + Decision Hub", GREEN),
    ("Phase 2", "4\uC8FC", "\uCC28\uD2B8 \uC2DC\uC2A4\uD15C + Alert Stream", BLUE),
    ("Phase 3", "2\uC8FC", "\uD3F4\uB9AC\uC2F1 + \uC0AC\uC6A9\uC790 \uD14C\uC2A4\uD2B8", CYAN),
]
for i, (phase, dur, desc, accent) in enumerate(timeline):
    tx = Inches(0.6 + i * 4.1)
    add_rounded_rect(slide, tx, Inches(6.6), Inches(3.8), Inches(0.5), BG_CARD, border_color=accent)
    add_text(slide, f"{phase}  {dur}", tx + Inches(0.15), Inches(6.63), Inches(1.5), Inches(0.2),
             size=10, color=accent, bold=True, font_name="Consolas")
    add_text(slide, desc, tx + Inches(1.5), Inches(6.63), Inches(2.2), Inches(0.2),
             size=10, color=WHITE)


# ============================================================
# SLIDE 18: CLOSING - NEXT STEPS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)

# Full green accent top
add_rect(slide, Inches(0), Inches(0), SW, Inches(0.06), GREEN)

# Title
add_text(slide, "NEXT STEPS", Inches(0.6), Inches(0.8), Inches(10), Inches(0.8),
         size=44, color=WHITE, bold=True, font_name="Segoe UI Light")
add_text(slide, "\uB2E4\uC74C \uB2E8\uACC4", Inches(0.6), Inches(1.5), Inches(10), Inches(0.5),
         size=18, color=GREEN, font_name="Segoe UI")

# 4 phases
next_phases = [
    ("PHASE 1", "\uC640\uC774\uC5B4\uD504\uB808\uC784 + Figma \uD504\uB85C\uD1A0\uD0C0\uC785", "2\uC8FC", GREEN),
    ("PHASE 2", "\uD575\uC2EC \uCEF4\uD3EC\uB10C\uD2B8 \uAD6C\uD604 (Decision Hub + Hero Chart)", "4\uC8FC", BLUE),
    ("PHASE 3", "\uC804\uCCB4 \uD1B5\uD569 + \uC0AC\uC6A9\uC790 \uD14C\uC2A4\uD2B8", "3\uC8FC", CYAN),
    ("PHASE 4", "\uD53C\uB4DC\uBC31 \uBC18\uC601 + \uC815\uC2DD \uB9B4\uB9AC\uC2A4", "1\uC8FC", YELLOW),
]

for i, (phase, desc, dur, accent) in enumerate(next_phases):
    cy = Inches(2.3 + i * 1.05)
    # Timeline dot
    add_circle(slide, Inches(1.0), cy + Inches(0.12), Inches(0.2), accent)
    add_text(slide, str(i+1), Inches(1.0), cy + Inches(0.12), Inches(0.2), Inches(0.2),
             size=10, color=BG_MAIN, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)
    # Connecting line
    if i < 3:
        add_rect(slide, Inches(1.08), cy + Inches(0.35), Pt(2), Inches(0.7), DIMMER)

    add_text(slide, phase, Inches(1.5), cy + Inches(0.05), Inches(1.5), Inches(0.3),
             size=14, color=accent, bold=True, font_name="Consolas")
    add_text(slide, desc, Inches(3.0), cy + Inches(0.05), Inches(7), Inches(0.3),
             size=13, color=WHITE)

    # Duration badge
    add_rounded_rect(slide, Inches(10.5), cy + Inches(0.05), Inches(1.2), Inches(0.35), BG_HOVER, border_color=accent)
    add_text(slide, dur, Inches(10.5), cy + Inches(0.08), Inches(1.2), Inches(0.3),
             size=11, color=accent, bold=True, font_name="Consolas", align=PP_ALIGN.CENTER)

# Total duration
add_rect(slide, Inches(0.6), Inches(6.0), Inches(12.1), Pt(1), GREEN)
add_text(slide, "TOTAL:  10\uC8FC  (2.5\uAC1C\uC6D4)", Inches(0.6), Inches(6.2), Inches(5), Inches(0.4),
         size=18, color=GREEN, bold=True, font_name="Consolas")

# Bottom bar
add_rect(slide, Inches(0), SH - Inches(0.8), SW, Inches(0.8), BG_CARD)
add_text(slide, "DESIGN STUDIO X", Inches(1.0), SH - Inches(0.65), Inches(3), Inches(0.4),
         size=14, color=GREEN, bold=True, font_name="Consolas")
add_text(slide, "\u00B7  Q-TRON QUANT  \u00B7  2026.04  \u00B7", Inches(4.0), SH - Inches(0.65), Inches(6), Inches(0.4),
         size=14, color=DIM, font_name="Consolas")
add_text(slide, "CONFIDENTIAL", Inches(10.0), SH - Inches(0.65), Inches(2.5), Inches(0.4),
         size=12, color=DIMMER, font_name="Consolas", align=PP_ALIGN.RIGHT)


# ── SAVE ──
output_path = r"C:\Q-TRON-32_ARCHIVE\kr-legacy\docs\Q-TRON_Design_Concept_Proposal.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
